import os.path
import pandas as pd
from pptx.util import Inches
from pptx import Presentation
from utils import PPT_Helper
import tqdm
import imageio
import cv2

prs = Presentation(r'C:/Users/cheny/Desktop/master_view.pptx')

_TGT = '//Salmon/User/Chen/Vessel_data/revised_nerve/crop_hipjiont_7.5cm_20220905'
os.makedirs(_TGT, exist_ok=True)
_PATIENT_LIST = '//Salmon/User/Chen/Vessel_data/revised_nerve/caseid_list_nerves.txt'
with open(_PATIENT_LIST) as f:
    case_IDs = f.read().splitlines()

for case_ID in tqdm.tqdm(case_IDs[0:20]):
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    PPT_Helper.add_title(slide, case_ID)

    # PPT_Helper.add_text(slide=slide, msg='vein_dice:{:.3f} '.format(vein_dice), left=20, top=-0.5, width=4, height=2,
    #                     font_size=16, is_bold=False)
    # PPT_Helper.add_text(slide=slide, msg='artery_dice:{:.3f} '.format(artery_dice), left=20, top=0.5, width=4, height=2,
    #
    #                     font_size=16, is_bold=False)

    pic_path = os.path.join(_TGT, 'visualization', case_ID + '_vessels_nerve.png')

    im = imageio.imread(pic_path)
    print(im.shape)

    ratio = im.shape[1] / im.shape[0]
    print(ratio)
    slide.shapes.add_picture(pic_path, Inches(4.00), Inches(0.18), Inches(6.48 * ratio),
                             Inches(6.48))

    movie_path = os.path.join(_TGT, 'movies', '{}.mp4'.format(case_ID))
    vidcap = cv2.VideoCapture(movie_path)
    success, image = vidcap.read()
    cv2.imwrite(os.path.join(_TGT, 'movies', '{}_1.png'.format(case_ID)), image)
    movie = slide.shapes.add_movie(movie_path, Inches(0.45), Inches(1.63), Inches(4.69), Inches(4.69),
                                   poster_frame_image=os.path.join(_TGT, 'movies', '{}_1.png'.format(case_ID)))
    # for shape in slide.placeholders:
    #     phf = shape.placeholder_format
    # print(f"{phf.idx}--{shape.name}--{phf.type}")
    # pic = slide.placeholders[1]
    # pic.insert_picture(os.path.join(_TGT,'visualization', case_ID + '_vessels_nerve.png'))
    # pic = slide.placeholders[2]
    # pic.insert_picture(os.path.join()
    # pic = slide.placeholders[12]
    # pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_vein_right_right_pelvis.png'))
    # pic = slide.placeholders[13]
    # pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '-vessels_label_lr_vein_right__right_pelvis.png'))
    # pic = slide.placeholders[14]
    # pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_artery_left_left_pelvis.png'))
    # pic = slide.placeholders[15]
    # pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '-vessels_label_lr_artery_left__left_pelvis.png'))
    # pic = slide.placeholders[16]
    # pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_artery_right_right_pelvis.png'))
    # pic = slide.placeholders[17]
    # pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '-vessels_label_lr_artery_right__right_pelvis.png'))
    prs.save(r'C:/Users/cheny/Desktop/nerve_test.pptx')
