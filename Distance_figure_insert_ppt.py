import os.path
import pandas as pd
from pptx import Presentation
from utils import PPT_Helper
import tqdm

def point2point_distance(p1, p2):
    distance = ((float(p1[0]) - float(p2[0])) ** 2 + (float(p1[1]) - float(p2[1])) ** 2 + (
                float(p1[2]) - float(p2[2])) ** 2) ** 0.5
    return distance


def readtxt(fpath):
    my_file = open(fpath, "r")
    data = my_file.read()
    data_into_list = data.replace('\n', ' ').split(",")
    print(data_into_list)
    my_file.close()
    return (data_into_list)


prs = Presentation(r'C:/Users/cheny/Desktop/Distance_template.pptx')

GT_distance_path = '//Salmon/User/Chen/Vessel_data/Dataset for Distance Assessment/crop_hipjiont_75mm_ver3/Polygons/PolyFigures'
GT_points_root = '//Salmon/User/Chen/Vessel_data/Dataset for Distance Assessment/crop_hipjiont_75mm_ver3/Polygons/PolyDistances'
# AUTO_data_root = '//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_5fold/seperation_left_right/Polygons/PolyFigures'
AUTO_data_root ='//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_5fold/seperation_left_right/Polygons/PolyFigures'
AUTO_points_root = '//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_5fold/seperation_left_right/Polygons/PolyDistances'
AUTO_data_muscles_root ='//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_muscles_5fold/seperation_left_right/Polygons/PolyFigures'
AUTO_points_muscles_root = '//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_muscles_5fold/seperation_left_right/Polygons/PolyDistances'
_PATIENT_LIST = '//Salmon/User/Chen/Vessel_data/revised_nerve/caseid_list_nerves.txt'
diff_csv_path = 'D:/temp/visualization'
Accuracy_path = '//Salmon/User/mazen/Segmentation/Codes/results/Nara/hip_vessels/18_5fold/seperation_left_right/Polygons/Analysis/Osaka_Nara_GT_Predicted_Mindistance.csv'
Dice = pd.read_csv(Accuracy_path, header=0, index_col=0)
# vein_left_GT = pd.read_csv(os.path.join(GT_distance_path, 'vein_left.csv'), header=0, index_col=0)
# artery_left_GT = pd.read_csv(os.path.join(GT_distance_path, 'artery_left.csv'), header=0, index_col=0)
# vein_right_GT = pd.read_csv(os.path.join(GT_distance_path, 'vein_right.csv'), header=0, index_col=0)
# artery_right_GT = pd.read_csv(os.path.join(GT_distance_path, 'artery_right.csv'), header=0, index_col=0)
# nerve_left_GT = pd.read_csv(os.path.join(GT_distance_path, 'nerve_left.csv'), header=0, index_col=0)
# nerve_right_GT = pd.read_csv(os.path.join(GT_distance_path, 'nerve_right.csv'), header=0, index_col=0)
#
# vein_left_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'vein_left.csv'), header=0, index_col=0)
# artery_left_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'artery_left.csv'), header=0, index_col=0)
# vein_right_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'vein_right.csv'), header=0, index_col=0)
# artery_right_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'artery_right.csv'), header=0, index_col=0)
# nerve_left_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'nerve_left.csv'), header=0, index_col=0)
# nerve_right_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'nerve_right.csv'), header=0, index_col=0)


with open(_PATIENT_LIST) as f:
    case_IDs = f.read().splitlines()
for case_ID in tqdm.tqdm(case_IDs):
    print (case_ID)
    case_ID = case_ID.replace('"','')
    # a= case_ID.replace('"','')
    # b=vein_left_GT.loc[a][0]
    # vein_left_GT.index[case_ID]
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    PPT_Helper.add_title(slide, case_ID)
    # a = vein_left_GT[vein_left_GT['ID'] == case_ID]
    # print('vein_dice:{:.2f} '.format(Dice.loc[case_ID.lower()][0]))
    # vein_dice = Dice.loc[case_ID.lower()][0]
    # artery_dice = Dice.loc[case_ID.lower()][1]
    # PPT_Helper.add_text(slide=slide, msg='vein_dice:{:.3f} '.format(vein_dice), left=20, top=-0.5, width=4, height=2,
    #                     font_size=16, is_bold=False)
    # PPT_Helper.add_text(slide=slide, msg='artery_dice:{:.3f} '.format(artery_dice), left=20, top=0.5, width=4, height=2,
    #                     font_size=16, is_bold=False)
    # print(vein_left_GT[case_ID][0])
    # a= point2point_distance(readtxt(os.path.join(GT_points_root,'{}_vein_left.txt'.format(case_ID))),
    #                          readtxt(os.path.join(AUTO_points_root,'{}_vein_left.txt'.format(case_ID))))
    # PPT_Helper.add_text(slide=slide, msg='Points Distance:{:.3f}mm '.format(
    #     point2point_distance(readtxt(os.path.join(GT_points_root, '{}_vein_left.txt'.format(case_ID))),
    #                          readtxt(os.path.join(AUTO_points_root, '{}_vein_left.txt'.format(case_ID))))),
    #                     left=8.5, top=10.7, width=4, height=2,
    #                     font_size=18, is_bold=True)
    # PPT_Helper.add_text(slide=slide, msg='Points_Distance:{:.3f}mm '.format(
    #     abs(vein_left_GT.loc[case_ID][0] - vein_left_Auto.loc[case_ID][0])),
    #                     left=8.5, top=10.7, width=4, height=2,
    #                     font_size=18, is_bold=True)
    # PPT_Helper.add_text(slide=slide, msg='Points Distance:{:.3f}mm '.format(
    #     point2point_distance(readtxt(os.path.join(GT_points_root, '{}_vein_right.txt'.format(case_ID))),
    #                          readtxt(os.path.join(AUTO_points_root, '{}_vein_right.txt'.format(case_ID))))),
    #                     left=24.5, top=10.7, width=4, height=2,
    #                     font_size=18, is_bold=True)
    # PPT_Helper.add_text(slide=slide, msg='Points Distance:{:.3f}mm '.format(
    #     point2point_distance(readtxt(os.path.join(GT_points_root, '{}_artery_left.txt'.format(case_ID))),
    #                          readtxt(os.path.join(AUTO_points_root, '{}_artery_left.txt'.format(case_ID))))),
    #                     left=8.5, top=18.7, width=4, height=2,
    #                     font_size=18, is_bold=True)
    # PPT_Helper.add_text(slide=slide, msg='Points Distance:{:.3f}mm '.format(
    #     point2point_distance(readtxt(os.path.join(GT_points_root, '{}_artery_right.txt'.format(case_ID))),
    #                          readtxt(os.path.join(AUTO_points_root, '{}_artery_right.txt'.format(case_ID))))),
    #                     left=24.5, top=18.7, width=4, height=2,
    #                     font_size=18, is_bold=True)
    PPT_Helper.add_text(slide=slide, msg='Points Distance:\n{:.3f}mm '.format(
        point2point_distance(readtxt(os.path.join(GT_points_root, '{}_nerve_left.txt'.format(case_ID))),
                             readtxt(os.path.join(AUTO_points_root, '{}_nerve_left.txt'.format(case_ID))))),
                        left=10.5, top=23.5, width=4, height=2,
                        font_size=18, is_bold=True)
    PPT_Helper.add_text(slide=slide, msg='Points Distance:\n{:.3f}mm '.format(
        point2point_distance(readtxt(os.path.join(GT_points_root, '{}_nerve_right.txt'.format(case_ID))),
                             readtxt(os.path.join(AUTO_points_root, '{}_nerve_right.txt'.format(case_ID))))),
                        left=33.5, top=23.5, width=4, height=2,
                        font_size=18, is_bold=True)
    PPT_Helper.add_text(slide=slide, msg='Points Distance:\n{:.3f}mm '.format(
        point2point_distance(readtxt(os.path.join(GT_points_root, '{}_nerve_left.txt'.format(case_ID))),
                             readtxt(os.path.join(AUTO_points_muscles_root, '{}_nerve_left.txt'.format(case_ID))))),
                        left=17.5, top=23.5, width=4, height=2,
                        font_size=18, is_bold=True)
    PPT_Helper.add_text(slide=slide, msg='Points Distance:\n{:.3f}mm '.format(
        point2point_distance(readtxt(os.path.join(GT_points_root, '{}_nerve_right.txt'.format(case_ID))),
                             readtxt(os.path.join(AUTO_points_muscles_root, '{}_nerve_right.txt'.format(case_ID))))),
                        left=39.5, top=23.5, width=4, height=2,
                        font_size=18, is_bold=True)
    for shape in slide.placeholders:
        phf = shape.placeholder_format
        print(f"{phf.idx}--{shape.name}--{phf.type}")
    pic = slide.placeholders[23]
    pic.insert_picture(os.path.join(GT_distance_path+'_nerve_back', case_ID + '_nerve_left_left_pelvis.png'))
    pic = slide.placeholders[47]
    pic.insert_picture(os.path.join(AUTO_data_root+'_nerve_back', case_ID + '_nerve_left_left_pelvis.png'))
    pic = slide.placeholders[48]
    pic.insert_picture(os.path.join(AUTO_data_muscles_root+'_nerve_back', case_ID + '_nerve_left_left_pelvis.png'))

    pic = slide.placeholders[49]
    pic.insert_picture(os.path.join(GT_distance_path+'_nerve_side', case_ID + '_nerve_left_left_pelvis.png'))
    pic = slide.placeholders[50]
    pic.insert_picture(os.path.join(AUTO_data_root+'_nerve_side', case_ID + '_nerve_left_left_pelvis.png'))
    pic = slide.placeholders[51]
    pic.insert_picture(os.path.join(AUTO_data_muscles_root+'_nerve_side', case_ID + '_nerve_left_left_pelvis.png'))

    pic = slide.placeholders[52]
    pic.insert_picture(os.path.join(GT_distance_path + '_nerve_back', case_ID + '_nerve_right_right_pelvis.png'))
    pic = slide.placeholders[53]
    pic.insert_picture(os.path.join(AUTO_data_root + '_nerve_back', case_ID + '_nerve_right_right_pelvis.png'))
    pic = slide.placeholders[54]
    pic.insert_picture(os.path.join(AUTO_data_muscles_root + '_nerve_back', case_ID + '_nerve_right_right_pelvis.png'))

    pic = slide.placeholders[55]
    pic.insert_picture(os.path.join(GT_distance_path+'_nerve_side', case_ID + '_nerve_right_right_pelvis.png'))
    pic = slide.placeholders[56]
    pic.insert_picture(os.path.join(AUTO_data_root+'_nerve_side', case_ID + '_nerve_right_right_pelvis.png'))
    pic = slide.placeholders[57]
    pic.insert_picture(os.path.join(AUTO_data_muscles_root+'_nerve_side', case_ID + '_nerve_right_right_pelvis.png'))
    prs.save(r'C:/Users/cheny/Desktop/Distance_Nerve_pelvis_Points_Distance_backside.pptx')
