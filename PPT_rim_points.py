import os.path
import numpy as np
import pandas as pd
from pptx import Presentation
from utils import PPT_Helper
import tqdm

def point2point_distance(p1, p2):
    distance = ((float(p1[0]) - float(p2[0])) ** 2 + (float(p1[1]) - float(p2[1])) ** 2 + (
                float(p1[2]) - float(p2[2])) ** 2) ** 0.5
    return distance


def readtxt(fpath):
    my_file = open(fpath, "r")
    data = my_file.read()
    data_into_list = data.replace('\n', ' ').split(",")
    print(data_into_list)
    my_file.close()
    return (data_into_list)


prs = Presentation(r'C:/Users/cheny/Desktop/Distance_template.pptx')

GT_distance_path = '//Salmon/User/Chen/Vessel_data/Dataset for Distance Assessment/crop_vessels_pelvis_36cases/seperation_left_right/Polygons/PolyFigures_CEll_20230220'
    # '//Salmon/User/Chen/Vessel_data/Nara_enhance_plain/Good alignment/crop_hip_joint/seperation_left_right/Polygons/PolyFigures_Risk'
    # '//Salmon/User/Chen/Vessel_data/Dataset for Distance Assessment/crop_vessels_pelvis_36cases/seperation_left_right/Polygons/PolyFigures_Risk'
# GT_points_root = '//Salmon/User/Chen/Vessel_data/revised_nerve/crop_hipjiont_75mm_ver3/Polygons/PolyDistances'
# AUTO_data_root = '//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_5fold/seperation_left_right/Polygons/PolyFigures'
# AUTO_data_root ='//Salmon/User/mazen/Segmentation/Codes/results/Nara/hip_vessels/18_5fold/seperation_left_right/Polygons/PolyFigures_Risk'
# AUTO_points_root = '//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_5fold/seperation_left_right/Polygons/PolyDistances'
# AUTO_data_muscles_root ='//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_muscles_5fold/seperation_left_right/Polygons/PolyFigures_nerve_back'
# AUTO_points_muscles_root = '//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_muscles_5fold/seperation_left_right/Polygons/PolyDistances'
# AUTO_data_root ='//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_5fold/seperation_left_right/Polygons/PolyFigures'
# AUTO_points_root = '//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_5fold/seperation_left_right/Polygons/PolyDistances'
Auto_distance_path = '//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/36_womuscles_cropped_150000/seperation_left_right/Polygons/PolyFigures_CEll_20230220'
    # '//Salmon/User/mazen/Segmentation/Codes/results/Nara/hip_vessels/18_5fold/seperation_left_right/Polygons/PolyFigures_Risk'
    # '//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/36_womuscles_cropped_150000/seperation_left_right/Polygons/PolyFigures_Risk'
_PATIENT_LIST = '//Salmon/User/Chen/Vessel_data/Dataset for Distance Assessment/crop_vessels_pelvis_36cases/caseid_list_vessels_36.txt'
    # '//Salmon/User/Chen/Vessel_data/Nara_enhance_plain/crop_cases_list.txt'
    # '//Salmon/User/Chen/Vessel_data/Dataset for Distance Assessment/crop_vessels_pelvis_36cases/caseid_list_vessels_36.txt'
Plots_path ='//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/36_womuscles_cropped_150000/seperation_left_right/Polygons/PolyDistances_20230214/Rim_distatnce'
    # '//Salmon/User/mazen/Segmentation/Codes/results/Nara/hip_vessels/18_5fold/seperation_left_right/Polygons/PolyDistances/GT_Predicted_Points_distance'
    # '//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/36_womuscles_cropped_150000/seperation_left_right/Polygons/PolyDistances/GT_Predicted_Points_distance'
diff_csv_path = 'D:/temp/visualization'
Accuracy_path = '//Salmon/User/Chen/Vessel_data/Dataset for Distance Assessment/crop_vessels_pelvis_36cases/seperation_left_right/Polygons/Analysis_20230214/Osaka_Nara_GT_Predicted_Mindistance.csv'
Dice = pd.read_csv(Accuracy_path, header=0, index_col=0)
Dice = Dice.set_index(['ID','vessel/side'])
replace_dict = {'N0018': '#2-1', 'N0024': '#2-2', 'N0047': '#2-3', 'N0056': '#2-4', 'N0074': '#2-5', 'N0076': '#2-6',
                  'N0091': '#2-7', 'N0094': '#2-8', 'N0107': '#2-9', 'N0108': '#2-10', 'N0116': '#2-11',
                  'N0132': '#2-12',
                  'N0133': '#2-13', 'N0140': '#2-14', 'N0144': '#2-15', 'N0171': '#2-16', 'N0180': '#2-17',
                  'N0187': '#2-18',
                  'k10387': '#1-1', 'k7510': '#1-2', 'k8559': '#1-3', 'k8574': '#1-4', 'k8699': '#1-5', 'k8748': '#1-6',
                  'k8772': '#1-7', 'k8895': '#1-8', 'k9020': '#1-9', 'k9089': '#1-10', 'k9162': '#1-11',
                  'k9193': '#1-12',
                  'k9204': '#1-13', 'k9622': '#1-14', 'k9831': '#1-15', 'k9861': '#1-16', 'k1565': '#1-17',
                  'k1585': '#1-18',
                  'k1631': '#1-19', 'k1657': '#1-20', 'k1665': '#1-21', 'k1677': '#1-22', 'k1712': '#1-23',
                  'k1756': '#1-24',
                  'k1796': '#1-25', 'k1802': '#1-26', 'k1870': '#1-27', 'k1873': '#1-28', 'k1647': '#1-29',
                  'k6940': '#1-30',
                  'k8041': '#1-31', 'k8454': '#1-32', 'k8795': '#1-33', 'k8892': '#1-34', 'k9086': '#1-35',
                  'k9339': '#1-36'}
# vein_left_GT = pd.read_csv(os.path.join(GT_distance_path, 'vein_left.csv'), header=0, index_col=0)
# artery_left_GT = pd.read_csv(os.path.join(GT_distance_path, 'artery_left.csv'), header=0, index_col=0)
# vein_right_GT = pd.read_csv(os.path.join(GT_distance_path, 'vein_right.csv'), header=0, index_col=0)
# artery_right_GT = pd.read_csv(os.path.join(GT_distance_path, 'artery_right.csv'), header=0, index_col=0)
# # nerve_left_GT = pd.read_csv(os.path.join(GT_distance_path, 'nerve_left.csv'), header=0, index_col=0)
# # nerve_right_GT = pd.read_csv(os.path.join(GT_distance_path, 'nerve_right.csv'), header=0, index_col=0)
#
# vein_left_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'vein_left.csv'), header=0, index_col=0)
# artery_left_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'artery_left.csv'), header=0, index_col=0)
# vein_right_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'vein_right.csv'), header=0, index_col=0)
# artery_right_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'artery_right.csv'), header=0, index_col=0)
# nerve_left_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'nerve_left.csv'), header=0, index_col=0)
# nerve_right_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'nerve_right.csv'), header=0, index_col=0)


with open(_PATIENT_LIST) as f:
    case_IDs = f.read().splitlines()
_cases_new_rule = np.array([replace_dict.get(item, item) for item in case_IDs])
# case_IDs = ['k10387']
for  i, case_ID in enumerate(tqdm.tqdm(case_IDs)):
    print (case_ID)
    case_ID = case_ID.replace('"','')
    ID_new = _cases_new_rule[i]
    # a= case_ID.replace('"','')
    # b=vein_left_GT.loc[a][0]
    # vein_left_GT.index[case_ID]
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    PPT_Helper.add_title(slide, ID_new)
    # a = vein_left_GT[vein_left_GT['ID'] == case_ID]
    # print('vein_dice:{:.2f} '.format(Dice.loc[case_ID.lower()][0]))
    # vein_dice = Dice.loc[case_ID.lower()][0]
    # artery_dice = Dice.loc[case_ID.lower()][1]
    # PPT_Helper.add_text(slide=slide, msg='vein_dice:{:.3f} '.format(vein_dice), left=20, top=-0.5, width=4, height=2,
    #                     font_size=16, is_bold=False)
    # PPT_Helper.add_text(slide=slide, msg='artery_dice:{:.3f} '.format(artery_dice), left=20, top=0.5, width=4, height=2,
    #                     font_size=16, is_bold=False)
    # print(vein_left_GT[case_ID][0])
    # a= point2point_distance(readtxt(os.path.join(GT_points_root,'{}_vein_left.txt'.format(case_ID))),
    #                          readtxt(os.path.join(AUTO_points_root,'{}_vein_left.txt'.format(case_ID))))
    # PPT_Helper.add_text(slide=slide, msg='Points Distance:{:.3f}mm '.format(
    #     point2point_distance(readtxt(os.path.join(GT_points_root, '{}_vein_left.txt'.format(case_ID))),
    #                          readtxt(os.path.join(AUTO_points_root, '{}_vein_left.txt'.format(case_ID))))),
    #                     left=8.5, top=10.7, width=4, height=2,
    #                     font_size=18, is_bold=True)
    # PPT_Helper.add_text(slide=slide, msg='Points_Distance:{:.3f}mm '.format(
    #     abs(vein_left_GT.loc[case_ID][0] - vein_left_Auto.loc[case_ID][0])),
    #                     left=8.5, top=10.7, width=4, height=2,
    #                     font_size=18, is_bold=True)
    # PPT_Helper.add_text(slide=slide, msg='Points Distance:{:.3f}mm '.format(
    #     point2point_distance(readtxt(os.path.join(GT_points_root, '{}_vein_right.txt'.format(case_ID))),
    #                          readtxt(os.path.join(AUTO_points_root, '{}_vein_right.txt'.format(case_ID))))),
    #                     left=24.5, top=10.7, width=4, height=2,
    #                     font_size=18, is_bold=True)
    dice = Dice.loc[(case_ID,'artery_left'),'Distance_error']
    PPT_Helper.add_text(slide=slide, msg='Error:{:.3f}mm '.format(
        Dice.loc[(case_ID, 'artery_left'), 'Distance_error']),
                        left=5.5, top=9.9, width=4, height=2,
                        font_size=18, is_bold=True)
    PPT_Helper.add_text(slide=slide, msg='Error:{:.3f}mm '.format(
        Dice.loc[(case_ID, 'artery_right'), 'Distance_error']),
                        left=26.5, top=9.9, width=4, height=2,
                        font_size=18, is_bold=True)
    PPT_Helper.add_text(slide=slide, msg='Error:{:.3f}mm '.format(
        Dice.loc[(case_ID, 'vein_left'), 'Distance_error']),
                        left=5.5, top=23.1, width=4, height=2,
                        font_size=18, is_bold=True)
    PPT_Helper.add_text(slide=slide, msg='Error:{:.3f}mm '.format(
        Dice.loc[(case_ID, 'vein_right'), 'Distance_error']),
                        left=26.5, top=23.1, width=4, height=2,
                        font_size=18, is_bold=True)
    for shape in slide.placeholders:
        phf = shape.placeholder_format
        # print(f"{phf.idx}--{shape.name}--{phf.type}")

    pic = slide.placeholders[22]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_artery_left_125_rim.png'))
    pic = slide.placeholders[23]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_artery_left_125_rim.png'))
    pic = slide.placeholders[17]
    pic.insert_picture(os.path.join(Plots_path, ID_new + '_artery_left__GT_Predicted_distance_dangerous.png'))
    pic = slide.placeholders[28]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_artery_right_125_rim.png'))
    pic = slide.placeholders[29]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_artery_right_125_rim.png'))
    pic = slide.placeholders[27]
    pic.insert_picture(os.path.join(Plots_path, ID_new + '_artery_right__GT_Predicted_distance_dangerous.png'))
    pic = slide.placeholders[25]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_vein_left_125_rim.png'))
    pic = slide.placeholders[26]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_vein_left_125_rim.png'))
    pic = slide.placeholders[24]
    pic.insert_picture(os.path.join(Plots_path, ID_new + '_vein_left__GT_Predicted_distance_dangerous.png'))
    pic = slide.placeholders[31]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_vein_right_125_rim.png'))
    pic = slide.placeholders[32]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_vein_right_125_rim.png'))
    pic = slide.placeholders[30]
    pic.insert_picture(os.path.join(Plots_path, ID_new + '_vein_left__GT_Predicted_distance_dangerous.png'))



    prs.save(r'//Salmon/User/Chen/Vessel_data/Dataset for Distance Assessment/crop_vessels_pelvis_36cases/seperation_left_right/Polygons/Analysis_20230214/Distance_Rim_Points_125_Osaka36_20230223.pptx')
