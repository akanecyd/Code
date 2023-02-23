import os.path
import pandas as pd
from pptx import Presentation
from utils import PPT_Helper
import tqdm

def point2point_distance(p1, p2):
    distance = ((float(p1[0]) - float(p2[0])) ** 2 + (float(p1[1]) - float(p2[1])) ** 2 + (
                float(p1[2]) - float(p2[2])) ** 2) ** 0.5
    return distance


def readtxt(fpath):
    my_file = open(fpath, "r")
    data = my_file.read()
    data_into_list = data.replace('\n', ' ').split(",")
    print(data_into_list)
    my_file.close()
    return (data_into_list)


prs = Presentation(r'C:/Users/cheny/Desktop/Distance_template.pptx')

GT_distance_path = '//Salmon/User/Chen/Vessel_data/Dataset for Distance Assessment/crop_vessels_pelvis_36cases/seperation_left_right/Polygons/PolyFigures_Risk'
    # '//Salmon/User/Chen/Vessel_data/Nara_enhance_plain/Good alignment/crop_hip_joint/seperation_left_right/Polygons/PolyFigures_Risk'
    # '//Salmon/User/Chen/Vessel_data/Dataset for Distance Assessment/crop_vessels_pelvis_36cases/seperation_left_right/Polygons/PolyFigures_Risk'
# GT_points_root = '//Salmon/User/Chen/Vessel_data/revised_nerve/crop_hipjiont_75mm_ver3/Polygons/PolyDistances'
# AUTO_data_root = '//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_5fold/seperation_left_right/Polygons/PolyFigures'
# AUTO_data_root ='//Salmon/User/mazen/Segmentation/Codes/results/Nara/hip_vessels/18_5fold/seperation_left_right/Polygons/PolyFigures_Risk'
# AUTO_points_root = '//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_5fold/seperation_left_right/Polygons/PolyDistances'
# AUTO_data_muscles_root ='//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_muscles_5fold/seperation_left_right/Polygons/PolyFigures_nerve_back'
# AUTO_points_muscles_root = '//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_muscles_5fold/seperation_left_right/Polygons/PolyDistances'
# AUTO_data_root ='//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_5fold/seperation_left_right/Polygons/PolyFigures'
# AUTO_points_root = '//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/20_w_rev_nerves_5fold/seperation_left_right/Polygons/PolyDistances'
Auto_distance_path = '//Salmon/User/mazen/Segmentation/Codes/results/Osaka/vessels/36_womuscles_cropped_150000/seperation_left_right/Polygons/PolyFigures_Risk'
    # '//Salmon/User/mazen/Segmentation/Codes/results/Nara/hip_vessels/18_5fold/seperation_left_right/Polygons/PolyFigures_Risk'
    # '//Salmon/User/mazen/Segmentation/Codes/results/Nara/hip_vessels/18_5fold/seperation_left_right/Polygons/PolyFigures_Risk'
_PATIENT_LIST = '//Salmon/User/Chen/Vessel_data/Dataset for Distance Assessment/crop_vessels_pelvis_36cases/caseid_list_vessels_36.txt'
    # '//Salmon/User/Chen/Vessel_data/Nara_enhance_plain/crop_cases_list.txt'
    # '//Salmon/User/Chen/Vessel_data/Dataset for Distance Assessment/crop_vessels_pelvis_36cases/caseid_list_vessels_36.txt'
diff_csv_path = 'D:/temp/visualization'
Accuracy_path = '//Salmon/User/mazen/Segmentation/Codes/results/Nara/hip_vessels/18_5fold/seperation_left_right/Polygons/Analysis/Osaka_Nara_GT_Predicted_Mindistance.csv'
Dice = pd.read_csv(Accuracy_path, header=0, index_col=0)
Dice = Dice.set_index(['ID','vessel/side'])
# vein_left_GT = pd.read_csv(os.path.join(GT_distance_path, 'vein_left.csv'), header=0, index_col=0)
# artery_left_GT = pd.read_csv(os.path.join(GT_distance_path, 'artery_left.csv'), header=0, index_col=0)
# vein_right_GT = pd.read_csv(os.path.join(GT_distance_path, 'vein_right.csv'), header=0, index_col=0)
# artery_right_GT = pd.read_csv(os.path.join(GT_distance_path, 'artery_right.csv'), header=0, index_col=0)
# # nerve_left_GT = pd.read_csv(os.path.join(GT_distance_path, 'nerve_left.csv'), header=0, index_col=0)
# # nerve_right_GT = pd.read_csv(os.path.join(GT_distance_path, 'nerve_right.csv'), header=0, index_col=0)
#
# vein_left_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'vein_left.csv'), header=0, index_col=0)
# artery_left_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'artery_left.csv'), header=0, index_col=0)
# vein_right_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'vein_right.csv'), header=0, index_col=0)
# artery_right_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'artery_right.csv'), header=0, index_col=0)
# nerve_left_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'nerve_left.csv'), header=0, index_col=0)
# nerve_right_Auto = pd.read_csv(os.path.join(Auto_distance_path, 'nerve_right.csv'), header=0, index_col=0)


with open(_PATIENT_LIST) as f:
    case_IDs = f.read().splitlines()
# case_IDs = ['k10387']
for case_ID in tqdm.tqdm(case_IDs):
    print (case_ID)
    case_ID = case_ID.replace('"','')
    # a= case_ID.replace('"','')
    # b=vein_left_GT.loc[a][0]
    # vein_left_GT.index[case_ID]
    slide = prs.slides.add_slide(prs.slide_layouts[12])
    PPT_Helper.add_title(slide, case_ID)
    # a = vein_left_GT[vein_left_GT['ID'] == case_ID]
    # print('vein_dice:{:.2f} '.format(Dice.loc[case_ID.lower()][0]))
    # vein_dice = Dice.loc[case_ID.lower()][0]
    # artery_dice = Dice.loc[case_ID.lower()][1]
    # PPT_Helper.add_text(slide=slide, msg='vein_dice:{:.3f} '.format(vein_dice), left=20, top=-0.5, width=4, height=2,
    #                     font_size=16, is_bold=False)
    # PPT_Helper.add_text(slide=slide, msg='artery_dice:{:.3f} '.format(artery_dice), left=20, top=0.5, width=4, height=2,
    #                     font_size=16, is_bold=False)
    # print(vein_left_GT[case_ID][0])
    # a= point2point_distance(readtxt(os.path.join(GT_points_root,'{}_vein_left.txt'.format(case_ID))),
    #                          readtxt(os.path.join(AUTO_points_root,'{}_vein_left.txt'.format(case_ID))))
    # PPT_Helper.add_text(slide=slide, msg='Points Distance:{:.3f}mm '.format(
    #     point2point_distance(readtxt(os.path.join(GT_points_root, '{}_vein_left.txt'.format(case_ID))),
    #                          readtxt(os.path.join(AUTO_points_root, '{}_vein_left.txt'.format(case_ID))))),
    #                     left=8.5, top=10.7, width=4, height=2,
    #                     font_size=18, is_bold=True)
    # PPT_Helper.add_text(slide=slide, msg='Points_Distance:{:.3f}mm '.format(
    #     abs(vein_left_GT.loc[case_ID][0] - vein_left_Auto.loc[case_ID][0])),
    #                     left=8.5, top=10.7, width=4, height=2,
    #                     font_size=18, is_bold=True)
    # PPT_Helper.add_text(slide=slide, msg='Points Distance:{:.3f}mm '.format(
    #     point2point_distance(readtxt(os.path.join(GT_points_root, '{}_vein_right.txt'.format(case_ID))),
    #                          readtxt(os.path.join(AUTO_points_root, '{}_vein_right.txt'.format(case_ID))))),
    #                     left=24.5, top=10.7, width=4, height=2,
    #                     font_size=18, is_bold=True)
    dice = Dice.loc[(case_ID,'artery_left'),'Distance_error']
    PPT_Helper.add_text(slide=slide, msg='Distance Error:{:.3f}mm '.format(
        Dice.loc[(case_ID, 'artery_left'), 'Distance_error']),
                        left=8.5, top=13.9, width=4, height=2,
                        font_size=18, is_bold=True)
    PPT_Helper.add_text(slide=slide, msg='Distance Error:{:.3f}mm '.format(
        Dice.loc[(case_ID, 'artery_right'), 'Distance_error']),
                        left=29.5, top=13.9, width=4, height=2,
                        font_size=18, is_bold=True)
    PPT_Helper.add_text(slide=slide, msg='Distance Error:{:.3f}mm '.format(
        Dice.loc[(case_ID, 'vein_left'), 'Distance_error']),
                        left=8.5, top=28.1, width=4, height=2,
                        font_size=18, is_bold=True)
    PPT_Helper.add_text(slide=slide, msg='Distance Error:{:.3f}mm '.format(
        Dice.loc[(case_ID, 'vein_right'), 'Distance_error']),
                        left=29.5, top=28.1, width=4, height=2,
                        font_size=18, is_bold=True)
    for shape in slide.placeholders:
        phf = shape.placeholder_format
        print(f"{phf.idx}--{shape.name}--{phf.type}")

    pic = slide.placeholders[10]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_artery_left_95_rim.png'))
    pic = slide.placeholders[11]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_artery_left_125_rim.png'))
    pic = slide.placeholders[12]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_artery_left_155_rim.png'))
    pic = slide.placeholders[13]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_artery_right_95_rim.png'))
    pic = slide.placeholders[14]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_artery_right_125_rim.png'))
    pic = slide.placeholders[15]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_artery_right_155_rim.png'))
    pic = slide.placeholders[16]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_artery_left_95_rim.png'))
    pic = slide.placeholders[17]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_artery_left_125_rim.png'))
    pic = slide.placeholders[18]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_artery_left_155_rim.png'))
    pic = slide.placeholders[19]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_artery_right_95_rim.png'))
    pic = slide.placeholders[20]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_artery_right_125_rim.png'))
    pic = slide.placeholders[21]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_artery_right_155_rim.png'))

    pic = slide.placeholders[22]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_vein_left_95_rim.png'))
    pic = slide.placeholders[23]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_vein_left_125_rim.png'))
    pic = slide.placeholders[24]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_vein_left_155_rim.png'))
    pic = slide.placeholders[25]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_vein_right_95_rim.png'))
    pic = slide.placeholders[26]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_vein_right_125_rim.png'))
    pic = slide.placeholders[27]
    pic.insert_picture(os.path.join(GT_distance_path, case_ID + '_vein_right_155_rim.png'))
    pic = slide.placeholders[28]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_vein_left_95_rim.png'))
    pic = slide.placeholders[29]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_vein_left_125_rim.png'))
    pic = slide.placeholders[30]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_vein_left_155_rim.png'))
    pic = slide.placeholders[31]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_vein_right_95_rim.png'))
    pic = slide.placeholders[32]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_vein_right_125_rim.png'))
    pic = slide.placeholders[33]
    pic.insert_picture(os.path.join(Auto_distance_path, case_ID + '_vein_right_155_rim.png'))

    prs.save(r'C:/Users/cheny/Desktop/Distance_vessel_rim_Distance_Osaka36.pptx')
